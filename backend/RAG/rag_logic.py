import os
import json
import csv
from typing import List, Dict, Any
import PyPDF2
from docx import Document
from pptx import Presentation
import openpyxl
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_community.docstore.document import Document as LangchainDocument
import logging

logger = logging.getLogger("omaya-rag")

class RAGManager:
    def __init__(self, storage_path: str = "backend/RAG/documents", index_path: str = "backend/RAG/faiss_index"):
        self.storage_path = storage_path
        self.index_path = index_path
        self.embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
        self.text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
        self.vector_store = None
        self.indexed_files = []

        if not os.path.exists(self.storage_path):
            os.makedirs(self.storage_path)

        self.load_index()

    def load_index(self):
        if os.path.exists(self.index_path):
            try:
                self.vector_store = FAISS.load_local(self.index_path, self.embeddings, allow_dangerous_deserialization=True)
                logger.info("RAG index loaded successfully")
            except Exception as e:
                logger.error(f"Error loading RAG index: {e}")
                self.vector_store = None
        else:
            logger.info("No RAG index found, starting fresh")

    def save_index(self):
        if self.vector_store:
            self.vector_store.save_local(self.index_path)
            logger.info("RAG index saved successfully")

    def process_file(self, file_path: str) -> List[LangchainDocument]:
        ext = os.path.splitext(file_path)[1].lower()
        text = ""

        try:
            if ext == ".pdf":
                with open(file_path, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        text += page.extract_text() + "\n"
            elif ext in [".doc", ".docx"]:
                doc = Document(file_path)
                for para in doc.paragraphs:
                    text += para.text + "\n"
            elif ext == ".pptx":
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            elif ext == ".xlsx":
                wb = openpyxl.load_workbook(file_path, data_only=True)
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    for row in ws.iter_rows(values_only=True):
                        text += " ".join([str(cell) for cell in row if cell is not None]) + "\n"
            elif ext == ".csv":
                with open(file_path, "r", encoding="utf-8") as f:
                    reader = csv.reader(f)
                    for row in reader:
                        text += " ".join(row) + "\n"
            elif ext == ".json":
                with open(file_path, "r", encoding="utf-8") as f:
                    data = json.load(f)
                    text = json.dumps(data, indent=2)
            elif ext == ".xml":
                with open(file_path, "r", encoding="utf-8") as f:
                    text = f.read()
            elif ext in [".txt", ".md"]:
                with open(file_path, "r", encoding="utf-8") as f:
                    text = f.read()
            else:
                logger.warning(f"Unsupported file format: {ext}")
                return []
        except Exception as e:
            logger.error(f"Error processing file {file_path}: {e}")
            return []

        if not text.strip():
            return []

        chunks = self.text_splitter.split_text(text)
        documents = [LangchainDocument(page_content=chunk, metadata={"source": os.path.basename(file_path)}) for chunk in chunks]
        return documents

    def index_files(self):
        all_documents = []
        new_files = []

        for filename in os.listdir(self.storage_path):
            file_path = os.path.join(self.storage_path, filename)
            if os.path.isfile(file_path) and filename not in self.indexed_files:
                docs = self.process_file(file_path)
                if docs:
                    all_documents.extend(docs)
                    new_files.append(filename)

        if all_documents:
            if self.vector_store is None:
                self.vector_store = FAISS.from_documents(all_documents, self.embeddings)
            else:
                self.vector_store.add_documents(all_documents)

            self.indexed_files.extend(new_files)
            self.save_index()
            return len(new_files), len(all_documents)
        return 0, 0

    def query(self, question: str, k: int = 4) -> List[Dict[str, Any]]:
        if not self.vector_store:
            return []

        results = self.vector_store.similarity_search_with_score(question, k=k)
        formatted_results = []
        for doc, score in results:
            formatted_results.append({
                "content": doc.page_content,
                "metadata": doc.metadata,
                "score": float(score)
            })
        return formatted_results

    def get_stats(self) -> Dict[str, Any]:
        num_files = len(os.listdir(self.storage_path))
        num_chunks = 0
        if self.vector_store:
            num_chunks = self.vector_store.index.ntotal

        return {
            "indexed_files": num_files,
            "total_chunks": num_chunks,
            "supported_formats": ["pdf", "doc", "docx", "pptx", "xlsx", "csv", "json", "xml", "txt", "md"]
        }

rag_manager = RAGManager()
